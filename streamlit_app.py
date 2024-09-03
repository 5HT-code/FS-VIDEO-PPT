import streamlit as st
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import moviepy.editor as mp
import tempfile

def save_slide_as_image(slide, output_path, size=(1280, 720)):
    """Convert slide to image."""
    # Create a blank image with white background
    img = Image.new('RGB', size, color='white')

    # Get the slide's background color (if any)
    if slide.background.fill.type == MSO_SHAPE_TYPE.FILL:
        bg_color = slide.background.fill.solid_fill_color.rgb
        img = Image.new('RGB', size, color=bg_color)

    # Save the blank image (this is where rendering shapes and images should be done)
    img.save(output_path)

def ppt_to_video_with_images(ppt_path, output_path, slide_duration, fade_duration):
    prs = Presentation(ppt_path)
    image_files = []

    with tempfile.TemporaryDirectory() as tmpdirname:
        for i, slide in enumerate(prs.slides):
            img_path = os.path.join(tmpdirname, f'slide_{i}.png')
            save_slide_as_image(slide, img_path)
            image_files.append(img_path)

        # Create video from images
        clips = []
        for img in image_files:
            clip = mp.ImageClip(img).set_duration(slide_duration)
            clips.append(clip.crossfadein(fade_duration).crossfadeout(fade_duration))

        video = mp.concatenate_videoclips(clips, padding=-fade_duration)
        video.write_videofile(output_path, fps=24)

def main():
    st.title("PPT to Video Converter")

    # File uploader for PPTX
    pptx_file = st.file_uploader("Upload PPTX file", type="pptx")

    # Parameters
    slide_duration = st.slider("Slide Duration (seconds)", 1, 10, 5)
    fade_duration = st.slider("Fade Duration (seconds)", 0.0, 2.0, 0.5)

    if st.button("Generate Video"):
        if pptx_file is None:
            st.error("Please upload a PPTX file.")
            return

        with st.spinner("Processing... This may take a while."):
            with tempfile.TemporaryDirectory() as tmpdirname:
                # Save uploaded PPTX
                pptx_path = os.path.join(tmpdirname, pptx_file.name)
                with open(pptx_path, "wb") as f:
                    f.write(pptx_file.getbuffer())

                # Convert PPT to video using images
                output_path = os.path.join(tmpdirname, "final_video.mp4")
                ppt_to_video_with_images(pptx_path, output_path, slide_duration, fade_duration)

                # Display video
                st.video(output_path)

                # Offer video for download
                with open(output_path, "rb") as video_file:
                    video_bytes = video_file.read()
                st.download_button(
                    label="Download Video",
                    data=video_bytes,
                    file_name=f"{os.path.splitext(pptx_file.name)[0]}_final.mp4",
                    mime="video/mp4"
                )

        st.success("Video generated successfully!")

if __name__ == "__main__":
    main()
